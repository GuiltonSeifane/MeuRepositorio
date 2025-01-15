from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import os

from pptx import Presentation
from datetime import datetime

class ZEROCLICKREPORT():
    """
    Module to collect data and create the presentation.
    """
    
    def __init__(self) -> None:
        
        self.presentation = Presentation()        
        self.presentation.slide_width = Inches(16)
        self.presentation.slide_height = Inches(9)
        
        #Creates a blank slide layout    
        self.title_slide_layout = self.presentation.slide_layouts[6]
        
        self.running_date = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        self.left = Inches(0)
        self.top = Inches(0) #+ Inches(4)
        self.width = Inches(16)
        self.height = Inches(0.8)
        
        self.footer_left = Inches(0)
        self.footer_top = Inches(8.4) #+ Inches(4)
        self.footer_width = Inches(7)
        self.footer_height = Inches(0.4)
                
        pass
    
    def create_slide(self, header="M-Pesa IT Operations Weekly Report", footer="M-PESA IT Operations Automated Weekly Report 2024", number=1):
        
        #Creates a slide
        self.slide = self.presentation.slides.add_slide(self.title_slide_layout)
        
        #Create a text box for slide title
        txBox_title = self.slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        
        #Fill in the color of the textbox (title)
        txBox_title.fill.solid()
        txBox_title.fill.fore_color.rgb = RGBColor(255, 0, 0)

        tf = txBox_title.text_frame

        p = tf.paragraphs[0]

        #Align Text to Center
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = header
        run.text.center
        
        title_font = run.font
        title_font.name = 'Calibri'
        title_font.size = Pt(36)
        title_font.bold = True
        title_font.italic = None  # cause value to be inherited from theme
        
        #Sets text color to white
        title_font.color.rgb = RGBColor(255, 255, 255)
        
        # Footer
        txBox_footer = self.slide.shapes.add_textbox(self.left, self.footer_top, self.footer_width, self.footer_height)
        
        #Fill in the color of the textbox (title)
        txBox_footer.fill.solid()
        txBox_footer.fill.fore_color.rgb = RGBColor(255, 255, 255)

        tpfotter = txBox_footer.text_frame

        pfotter = tpfotter.paragraphs[0]

        #Align Text to Center
        pfotter.alignment = PP_ALIGN.CENTER

        runfotter = pfotter.add_run()
        runfotter.text = f'{footer} - page {number}'
        runfotter.text.center
        
        #Add image footer
        self.add_image(left=13.5, top=7.8,height=2)
        
        return f"Slide for {header} report created."

    def add_image(self, image="./images/vm_mpesa_logos.png", left = 0.75, top = 1.25, height=0.5, width=0):
        
        #Add Image to the slide
        # pic = self.slide.shapes.add_picture(image.format('SPLIT'), Inches(left), Inches(top), Inches(height))
        
        # return
        if width == 0:
            
            picture = self.slide.shapes.add_picture(image, Inches(left), Inches(top), Inches(height))
        
        else:
            
            picture = self.slide.shapes.add_picture(image, Inches(left), Inches(top), Inches(height), Inches(width))
            # picture = self.slide.shapes.add_picture('system_availability.jpg', Inches(1), Inches(1))
        
    def add_excel_file(self, excel_file="abc.xlsx", left = 0.75, top = 1.25, height=0.5, width=14):
        
        self.slide.shapes.add_ole_object('datasource/M-PESA E2E Service Availability.xls', 'Prog_ID', Inches(left), Inches(top), Inches(height), Inches(width))
        
    def create_header(self):
        
        return "Header Created"

    def create_footer(self):
        
        return "Footer Created"
    
    def save_presentation(self, name=os.path.join("presentations", "M-PESA-IT-OPS_Report")):
        
        self.presentation.save(f"{name}_{self.running_date}.pptx")
        
        return "Presentation Saved"