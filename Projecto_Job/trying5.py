import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime

# Função para processar e limpar os dados do arquivo Excel
def process_excel(file_path):
    df = pd.read_excel(file_path)
    # Remover duplicados
    df = df.drop_duplicates()

    # Preencher valores ausentes com base no tipo de dado
    for column in df.columns:
        if pd.api.types.is_numeric_dtype(df[column]):
            df[column].fillna(0, inplace=True)  # Preencher valores numéricos com 0
        elif pd.api.types.is_datetime64_any_dtype(df[column]):
            df[column].fillna(pd.NaT, inplace=True)  # Preencher valores datetime com NaT
        elif pd.api.types.is_timedelta64_dtype(df[column]):
            df[column].fillna(pd.Timedelta(0), inplace=True)  # Preencher timedelta com 0
        else:
            df[column].fillna("N/A", inplace=True)  # Preencher outros tipos com "N/A"

    return df

# Classe para gerar a apresentação PowerPoint
class ZEROCLICKREPORT:
    def __init__(self) -> None:
        self.presentation = Presentation()        
        self.presentation.slide_width = Inches(16)
        self.presentation.slide_height = Inches(9)
        
        # Layout de slide vazio
        self.title_slide_layout = self.presentation.slide_layouts[6]
        self.running_date = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Definir posições e tamanhos de elementos
        self.left = Inches(0)
        self.top = Inches(0)
        self.width = Inches(16)
        self.height = Inches(0.8)
        self.footer_left = Inches(0)
        self.footer_top = Inches(8.4)
        self.footer_width = Inches(7)
        self.footer_height = Inches(0.4)

    # Função para criar o slide com cabeçalho e rodapé
    def create_slide(self, header="M-Pesa IT Operations Weekly Report", footer="M-PESA IT Operations Automated Weekly Report 2024", number=1):
        self.slide = self.presentation.slides.add_slide(self.title_slide_layout)

        # Adicionar caixa de texto para o título
        txBox_title = self.slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        txBox_title.fill.solid()
        txBox_title.fill.fore_color.rgb = RGBColor(255, 0, 0)

        tf = txBox_title.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.name = 'Calibri'
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Adicionar rodapé
        txBox_footer = self.slide.shapes.add_textbox(self.left, self.footer_top, self.footer_width, self.footer_height)
        txBox_footer.fill.solid()
        txBox_footer.fill.fore_color.rgb = RGBColor(255, 255, 255)

        tpfotter = txBox_footer.text_frame
        pfotter = tpfotter.paragraphs[0]
        pfotter.alignment = PP_ALIGN.CENTER
        runfotter = pfotter.add_run()
        runfotter.text = f'{footer} - page {number}'
        
        self.add_image(left=13.5, top=7.8, height=2)

    # Função para adicionar imagens aos slides
    def add_image(self, image="./images/vm_mpesa_logos.png", left=0.75, top=1.25, height=0.5, width=0):
        if width == 0:
            self.slide.shapes.add_picture(image, Inches(left), Inches(top), Inches(height))
        else:
            self.slide.shapes.add_picture(image, Inches(left), Inches(top), Inches(height), Inches(width))

    # Função para adicionar conteúdo de uma linha do Excel ao slide
    def add_row_content(self, row_data, left=1, top=2, width=14, height=6):
        txBox_content = self.slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox_content.text_frame
        p = tf.paragraphs[0]
        p.text = row_data
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Função para salvar a apresentação
    def save_presentation(self, name=os.path.join("presentations", "M-PESA-IT-OPS_Report")):
        self.presentation.save(f"{name}_{self.running_date}.pptx")

# Função principal
def main():
    try:
        # Caminho para o arquivo Excel
        excel_file = r"C:\Users\seifg001\Downloads\List of incidents_Jan.xlsx"

        # Processar o arquivo Excel
        df = process_excel(excel_file)

        # Gerar relatório PowerPoint
        zeroClickPresentation = ZEROCLICKREPORT()

        # Criar um slide para cada linha do DataFrame
        for index, row in df.iterrows():
            # Criar um slide com o conteúdo da linha
            slide_title = f"Linha {index + 1}"
            slide_content = "\n".join([f"{col}: {row[col]}" for col in df.columns])  # Formatar conteúdo da linha
            zeroClickPresentation.create_slide(header=slide_title, footer="Relatório Automático", number=index + 1)
            
            # Adicionar o conteúdo da linha ao slide
            zeroClickPresentation.add_row_content(slide_content)

        # Salvar apresentação PowerPoint
        zeroClickPresentation.save_presentation()

        print("Relatório PowerPoint gerado com sucesso!")

    except Exception as e:
        print(f"Erro: {e}")

if __name__ == "__main__":
    main()