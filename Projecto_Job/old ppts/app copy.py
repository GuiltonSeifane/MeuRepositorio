import collections
import collections.abc
from pptx import Presentation

presentation = Presentation()
slide_layout = presentation.slide_layouts[3]
slide = presentation.slides.add_slide(slide_layout)
content1 = slide.placeholders[1]
content1.text = 'First content'
content2 = slide.placeholders[2]
content2.text = 'Second content'
presentation.save('presentation.pptx')
print("Presentation created successfully.")